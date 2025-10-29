#!/usr/bin/env python3
"""
Convert markdown file to PowerPoint presentation
Usage: python md_to_pptx.py <input.md> [notes.md]
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import re
import sys
import os

def parse_markdown_slides(md_file):
    """Parse markdown file and extract slides"""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by slide separator
    slides = content.split('---')

    # Remove the frontmatter (first section with marp config)
    slides = [s.strip() for s in slides if s.strip()]
    if slides and 'marp:' in slides[0]:
        slides = slides[1:]

    parsed_slides = []
    for slide in slides:
        if not slide.strip():
            continue

        lines = slide.strip().split('\n')

        # Extract title (lines starting with # or ##)
        title = ""
        content_lines = []

        for line in lines:
            if line.startswith('# '):
                title = line.lstrip('#').strip()
            elif line.startswith('## '):
                title = line.lstrip('#').strip()
            elif line.startswith('### '):
                # Subtitle/section header
                content_lines.append(line.lstrip('#').strip())
            else:
                content_lines.append(line)

        # Join content and clean up
        content = '\n'.join(content_lines).strip()

        parsed_slides.append({
            'title': title,
            'content': content
        })

    return parsed_slides

def parse_markdown_notes(notes_file):
    """Parse markdown notes file and extract notes for each slide"""
    if not notes_file or not os.path.exists(notes_file):
        return []

    with open(notes_file, 'r', encoding='utf-8') as f:
        content = f.read()

    # Split by slide separator
    notes_sections = content.split('---')

    # Remove the frontmatter if present
    notes_sections = [s.strip() for s in notes_sections if s.strip()]
    if notes_sections and 'marp:' in notes_sections[0]:
        notes_sections = notes_sections[1:]

    parsed_notes = []
    for section in notes_sections:
        if not section.strip():
            continue

        # Remove the title line (first line starting with ##)
        lines = section.strip().split('\n')
        note_lines = []

        for line in lines:
            # Skip the slide title heading
            if line.startswith('## Slide'):
                continue
            elif line.startswith('# '):
                continue
            elif line.startswith('## '):
                continue
            else:
                note_lines.append(line)

        # Join note lines and clean up
        note_text = '\n'.join(note_lines).strip()
        parsed_notes.append(note_text)

    return parsed_notes

def create_presentation(slides, output_file, notes=None):
    """Create PowerPoint presentation from parsed slides"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    for idx, slide_data in enumerate(slides):
        # Use title and content layout
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)

        # Set title
        title = slide.shapes.title
        title.text = slide_data['title']
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True

        # Set content
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()

        # Parse content and add formatted text
        content = slide_data['content']

        # Split content into lines
        lines = content.split('\n')

        for i, line in enumerate(lines):
            if not line.strip():
                continue

            # Add paragraph
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            # Handle bullet points
            if line.strip().startswith('- '):
                text = line.strip()[2:]
                p.text = text
                p.level = 0
                p.font.size = Pt(18)
            elif line.strip().startswith('  - '):
                text = line.strip()[2:]
                p.text = text
                p.level = 1
                p.font.size = Pt(16)
            # Handle numbered lists
            elif re.match(r'^\d+\.', line.strip()):
                text = re.sub(r'^\d+\.\s*', '', line.strip())
                p.text = text
                p.level = 0
                p.font.size = Pt(18)
            # Handle bold text (markdown **text**)
            elif line.strip().startswith('**') and line.strip().endswith('**'):
                text = line.strip().strip('*')
                p.text = text
                p.font.size = Pt(20)
                p.font.bold = True
            else:
                p.text = line.strip()
                p.font.size = Pt(18)

        # Add speaker notes if available
        if notes and idx < len(notes) and notes[idx]:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes[idx]

    # Save presentation
    prs.save(output_file)
    print(f"PowerPoint presentation saved to: {output_file}")

def main():
    # Check for command-line argument
    if len(sys.argv) < 2:
        print("Usage: python md_to_pptx.py <slides.md> [notes.md]")
        print("Example: python md_to_pptx.py llm-simple.md llm-simple-notes.md")
        sys.exit(1)

    input_file = sys.argv[1]
    notes_file = sys.argv[2] if len(sys.argv) > 2 else None

    # Check if input file exists
    if not os.path.exists(input_file):
        print(f"Error: File '{input_file}' not found")
        sys.exit(1)

    # Check if notes file exists (if provided)
    if notes_file and not os.path.exists(notes_file):
        print(f"Warning: Notes file '{notes_file}' not found. Continuing without notes.")
        notes_file = None

    # Generate output filename by replacing .md with .pptx
    if input_file.endswith('.md'):
        output_file = input_file[:-3] + '.pptx'
    else:
        output_file = input_file + '.pptx'

    print(f"Reading {input_file}...")
    slides = parse_markdown_slides(input_file)
    print(f"Found {len(slides)} slides")

    notes = []
    if notes_file:
        print(f"Reading notes from {notes_file}...")
        notes = parse_markdown_notes(notes_file)
        print(f"Found {len(notes)} note sections")

    print(f"Creating PowerPoint presentation...")
    create_presentation(slides, output_file, notes)
    print("Done!")

if __name__ == '__main__':
    main()
